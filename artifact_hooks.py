import builtins
import os
import sys
import time
from typing import Optional

import theme


_WORKDIR: Optional[str] = None
_ARTIFACT_COUNTER = 0
_AUTO_HOOKS_INSTALLED = False
_ORIG_IMPORT = builtins.__import__
_REPORTLAB_TTF_REGISTERED = False


def _apply_matplotlib_defaults() -> None:
    try:
        import matplotlib as mpl
    except Exception:
        return

    try:
        try:
            mpl.rcdefaults()
        except Exception:
            pass

        try:
            from cycler import cycler

            mpl.rcParams["axes.prop_cycle"] = cycler(color=theme.CHART_COLORWAY_HEX)
        except Exception:
            pass

        mpl.rcParams.update(
            {
                "font.family": "sans-serif",
                "font.sans-serif": theme.TEXT_FONT_STACK,
                "figure.dpi": 160,
                "savefig.dpi": 160,
                "figure.facecolor": theme.COLOR_BACKGROUND_HEX,
                "axes.facecolor": theme.COLOR_BACKGROUND_HEX,
                "axes.edgecolor": theme.COLOR_BORDER_HEX,
                "axes.labelcolor": theme.COLOR_TEXT_HEX,
                "xtick.color": theme.COLOR_MUTED_HEX,
                "ytick.color": theme.COLOR_MUTED_HEX,
                "text.color": theme.COLOR_TEXT_HEX,
                "axes.spines.top": False,
                "axes.spines.right": False,
                "axes.spines.left": False,
                "axes.spines.bottom": False,
                "axes.grid": True,
                "grid.color": theme.COLOR_GRID_HEX,
                "grid.linewidth": 1.0,
                "grid.alpha": 1.0,
                "axes.axisbelow": True,
                "lines.linewidth": 2.0,
                "lines.solid_capstyle": "round",
                "lines.solid_joinstyle": "round",
                "lines.antialiased": True,
                "patch.antialiased": True,
                "axes.titlesize": 14,
                "axes.titleweight": "semibold",
                "axes.labelsize": 11,
                "xtick.labelsize": 10,
                "ytick.labelsize": 10,
                "legend.frameon": False,
            }
        )
    except Exception:
        return

    try:
        import matplotlib.axes as maxes
        import matplotlib.figure as mfigure
    except Exception:
        return

    if not getattr(maxes.Axes.set_title, "__croki_patched__", False):
        orig_set_title = maxes.Axes.set_title

        def set_title(self, label, fontdict=None, loc=None, pad=None, y=None, **kwargs):
            if fontdict is None:
                fontdict = {}
            if (
                "fontfamily" not in kwargs
                and "family" not in kwargs
                and "fontname" not in kwargs
                and "fontfamily" not in fontdict
                and "family" not in fontdict
                and "fontname" not in fontdict
            ):
                fontdict = dict(fontdict)
                fontdict["fontfamily"] = theme.HEADING_FONT_FAMILY
            return orig_set_title(self, label, fontdict=fontdict, loc=loc, pad=pad, y=y, **kwargs)

        set_title.__croki_patched__ = True
        maxes.Axes.set_title = set_title

    if not getattr(mfigure.Figure.suptitle, "__croki_patched__", False):
        orig_suptitle = mfigure.Figure.suptitle

        def suptitle(self, t, **kwargs):
            if "fontfamily" not in kwargs and "family" not in kwargs and "fontname" not in kwargs:
                kwargs = dict(kwargs)
                kwargs["fontfamily"] = theme.HEADING_FONT_FAMILY
            return orig_suptitle(self, t, **kwargs)

        suptitle.__croki_patched__ = True
        mfigure.Figure.suptitle = suptitle


def _apply_matplotlib_figure_defaults(fig) -> None:
    try:
        fig.set_facecolor(theme.COLOR_BACKGROUND_HEX)
        for ax in getattr(fig, "axes", []):
            try:
                ax.set_facecolor(theme.COLOR_BACKGROUND_HEX)
            except Exception:
                pass

            try:
                for spine in ax.spines.values():
                    spine.set_visible(False)
            except Exception:
                pass

            try:
                ax.set_axisbelow(True)
            except Exception:
                pass

            try:
                ax.grid(True, axis="y", color=theme.COLOR_GRID_HEX, linewidth=1.0)
                ax.grid(False, axis="x")
            except Exception:
                pass

            try:
                ax.tick_params(axis="both", which="both", length=0, colors=theme.COLOR_MUTED_HEX)
            except Exception:
                pass

            try:
                ax.margins(x=0.02)
            except Exception:
                pass
    except Exception:
        return


def _apply_plotly_defaults() -> None:
    try:
        import plotly.io as pio
    except Exception:
        return

    try:
        tmpl = pio.templates["plotly_white"]
        tmpl.layout.update(
            font=dict(
                family=", ".join(theme.TEXT_FONT_STACK),
                size=13,
                color=theme.COLOR_TEXT_HEX,
            ),
            title=dict(
                font=dict(
                    family=", ".join(theme.HEADING_FONT_STACK),
                )
            ),
            paper_bgcolor=theme.COLOR_BACKGROUND_HEX,
            plot_bgcolor=theme.COLOR_BACKGROUND_HEX,
            colorway=theme.CHART_COLORWAY_HEX,
            margin=dict(l=56, r=24, t=56, b=48),
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
        )
        tmpl.layout.xaxis.update(
            showgrid=True,
            gridcolor=theme.COLOR_GRID_HEX,
            gridwidth=1,
            zeroline=False,
            showline=False,
            ticks="",
            tickfont=dict(color=theme.COLOR_MUTED_HEX),
        )
        tmpl.layout.yaxis.update(
            showgrid=True,
            gridcolor=theme.COLOR_GRID_HEX,
            gridwidth=1,
            zeroline=False,
            showline=False,
            ticks="",
            tickfont=dict(color=theme.COLOR_MUTED_HEX),
        )
        pio.templates.default = "plotly_white"
    except Exception:
        return


def _apply_plotly_figure_defaults(fig) -> None:
    try:
        family_text = ", ".join(theme.TEXT_FONT_STACK)
        family_heading = ", ".join(theme.HEADING_FONT_STACK)

        font_family = None
        try:
            font_family = fig.layout.font.family
        except Exception:
            font_family = None

        title_family = None
        try:
            title_family = fig.layout.title.font.family
        except Exception:
            title_family = None

        patch: dict = {}
        try:
            if getattr(fig.layout, "template", None) is None:
                patch["template"] = "plotly_white"
        except Exception:
            patch["template"] = "plotly_white"
        if not font_family:
            patch["font"] = dict(family=family_text, size=14, color=theme.COLOR_TEXT_HEX)
        if not title_family:
            patch["title_font"] = dict(family=family_heading)

        try:
            if getattr(fig.layout, "paper_bgcolor", None) is None:
                patch["paper_bgcolor"] = theme.COLOR_BACKGROUND_HEX
        except Exception:
            patch["paper_bgcolor"] = theme.COLOR_BACKGROUND_HEX

        try:
            if getattr(fig.layout, "plot_bgcolor", None) is None:
                patch["plot_bgcolor"] = theme.COLOR_BACKGROUND_HEX
        except Exception:
            patch["plot_bgcolor"] = theme.COLOR_BACKGROUND_HEX

        if patch:
            fig.update_layout(**patch)
    except Exception:
        return


def _patch_reportlab():
    try:
        from reportlab.pdfgen.canvas import Canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except Exception:
        return

    if getattr(Canvas.__init__, "__croki_patched__", False):
        return

    def _register_ttf(font_name: str, path: str) -> None:
        try:
            pdfmetrics.getFont(font_name)
            return
        except Exception:
            pass
        try:
            pdfmetrics.registerFont(TTFont(font_name, path))
        except Exception:
            return

    orig_init = Canvas.__init__

    def __init__(self, *args, **kwargs):
        global _REPORTLAB_TTF_REGISTERED
        if not _REPORTLAB_TTF_REGISTERED:
            _register_ttf("Lato", "/usr/share/fonts/truetype/lato/Lato-Regular.ttf")
            _register_ttf("Unbounded", "/usr/local/share/fonts/unbounded/Unbounded-Regular.ttf")
            try:
                pdfmetrics.getFont("Lato")
                _REPORTLAB_TTF_REGISTERED = True
            except Exception:
                _REPORTLAB_TTF_REGISTERED = False
        orig_init(self, *args, **kwargs)
        try:
            r, g, b = theme.hex_to_rgb_float(theme.COLOR_TEXT_HEX)
            self.setFillColorRGB(r, g, b)
        except Exception:
            pass
        try:
            self.setFont(theme.TEXT_FONT_FAMILY, 11)
        except Exception:
            try:
                self.setFont("Helvetica", 11)
            except Exception:
                pass

    __init__.__croki_patched__ = True
    Canvas.__init__ = __init__


def _patch_docx():
    try:
        import docx
        import docx.api as api
        from docx.shared import RGBColor
    except Exception:
        return

    if getattr(api.Document, "__croki_patched__", False):
        return

    text_rgb = None
    try:
        text_rgb = RGBColor(*theme.hex_to_rgb_int(theme.COLOR_TEXT_HEX))
    except Exception:
        text_rgb = None

    orig_document = api.Document

    def Document(*args, **kwargs):
        doc = orig_document(*args, **kwargs)
        try:
            styles = doc.styles
            normal = styles["Normal"]
            normal.font.name = theme.TEXT_FONT_FAMILY
            if text_rgb is not None:
                normal.font.color.rgb = text_rgb
        except Exception:
            pass

        for i in range(1, 7):
            try:
                st = doc.styles[f"Heading {i}"]
                st.font.name = theme.HEADING_FONT_FAMILY
                if text_rgb is not None:
                    st.font.color.rgb = text_rgb
            except Exception:
                continue
        return doc

    Document.__croki_patched__ = True
    api.Document = Document
    docx.Document = Document


def _patch_pptx():
    try:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import PP_PLACEHOLDER
        from pptx.presentation import Presentation as PresentationClass
    except Exception:
        return

    if getattr(PresentationClass.save, "__croki_patched__", False):
        return

    text_rgb = None
    try:
        text_rgb = RGBColor(*theme.hex_to_rgb_int(theme.COLOR_TEXT_HEX))
    except Exception:
        text_rgb = None

    def _is_title_shape(shape) -> bool:
        try:
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                return True
        except Exception:
            pass
        return False

    def _apply_presentation_defaults(prs) -> None:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                tf = shape.text_frame
                is_title = _is_title_shape(shape)
                for p in tf.paragraphs:
                    try:
                        f = p.font
                        if not f.name:
                            f.name = theme.HEADING_FONT_FAMILY if is_title else theme.TEXT_FONT_FAMILY
                        if text_rgb is not None and (not f.color or not f.color.rgb):
                            f.color.rgb = text_rgb
                    except Exception:
                        pass

                    for r in p.runs:
                        try:
                            rf = r.font
                            if not rf.name:
                                rf.name = theme.HEADING_FONT_FAMILY if is_title else theme.TEXT_FONT_FAMILY
                            if text_rgb is not None and (not rf.color or not rf.color.rgb):
                                rf.color.rgb = text_rgb
                        except Exception:
                            continue

    orig_save = PresentationClass.save

    def save(self, *args, **kwargs):
        try:
            _apply_presentation_defaults(self)
        except Exception:
            pass
        return orig_save(self, *args, **kwargs)

    save.__croki_patched__ = True
    PresentationClass.save = save


def _ensure_matplotlib_env(workdir: str) -> None:
    if not os.environ.get("MPLBACKEND"):
        os.environ["MPLBACKEND"] = "Agg"
    if not os.environ.get("HOME"):
        os.environ["HOME"] = workdir

    xdg_config = os.environ.get("XDG_CONFIG_HOME")
    if not xdg_config:
        xdg_config = os.path.join(workdir, ".config")
        os.environ["XDG_CONFIG_HOME"] = xdg_config
    os.makedirs(xdg_config, exist_ok=True)

    mplconfig = os.environ.get("MPLCONFIGDIR")
    if not mplconfig:
        mplconfig = os.path.join(workdir, ".mplconfig")
        os.environ["MPLCONFIGDIR"] = mplconfig
    os.makedirs(mplconfig, exist_ok=True)

    xdg_cache = os.environ.get("XDG_CACHE_HOME")
    if not xdg_cache:
        xdg_cache = os.path.join(workdir, ".cache")
        os.environ["XDG_CACHE_HOME"] = xdg_cache
    os.makedirs(xdg_cache, exist_ok=True)


def _new_artifact_path(prefix: str, ext: str) -> str:
    global _ARTIFACT_COUNTER
    _ARTIFACT_COUNTER += 1
    workdir = _WORKDIR or "/workspace"
    name = f"{prefix}_{int(time.time() * 1000)}_{_ARTIFACT_COUNTER}.{ext}"
    return os.path.join(workdir, name)


def _patch_matplotlib():
    _ensure_matplotlib_env(_WORKDIR or "/workspace")
    try:
        import matplotlib

        try:
            matplotlib.use("Agg", force=True)
        except Exception:
            pass
        import matplotlib.pyplot as plt
    except Exception:
        return

    if getattr(plt.show, "__croki_patched__", False):
        return

    _apply_matplotlib_defaults()

    orig_show = plt.show

    def show(*args, **kwargs):
        nums = list(plt.get_fignums())
        if not nums:
            return orig_show(*args, **kwargs)
        for num in nums:
            _apply_matplotlib_defaults()
            fig = plt.figure(num)
            _apply_matplotlib_figure_defaults(fig)
            path = _new_artifact_path("plot", "png")
            try:
                fig.savefig(path)
            except Exception:
                _ensure_matplotlib_env(_WORKDIR or "/workspace")
                fig.savefig(path)
        plt.close("all")
        return None

    show.__croki_patched__ = True
    plt.show = show


def _patch_pil():
    try:
        from PIL import Image
    except Exception:
        return

    if getattr(Image.Image.show, "__croki_patched__", False):
        return

    orig_show = Image.Image.show

    def show(self, *args, **kwargs):
        path = _new_artifact_path("image", "png")
        try:
            self.save(path)
            return None
        except Exception:
            try:
                self.convert("RGB").save(path)
                return None
            except Exception:
                return orig_show(self, *args, **kwargs)

    show.__croki_patched__ = True
    Image.Image.show = show


def _patch_plotly():
    try:
        import plotly.basedatatypes as basedatatypes
    except Exception:
        return

    if getattr(basedatatypes.BaseFigure.show, "__croki_patched__", False):
        return

    _apply_plotly_defaults()

    orig_show = basedatatypes.BaseFigure.show

    def show(self, *args, **kwargs):
        path = _new_artifact_path("plotly", "html")
        try:
            _apply_plotly_figure_defaults(self)
            self.write_html(path, include_plotlyjs="cdn", full_html=True)
            return None
        except Exception:
            return orig_show(self, *args, **kwargs)

    show.__croki_patched__ = True
    basedatatypes.BaseFigure.show = show


def _croki_import(name, globals=None, locals=None, fromlist=(), level=0):
    mod = _ORIG_IMPORT(name, globals, locals, fromlist, level)
    try:
        if name == "matplotlib.pyplot" or (name == "matplotlib" and "pyplot" in (fromlist or ())):
            _patch_matplotlib()
        if name.startswith("PIL"):
            _patch_pil()
        if name.startswith("plotly"):
            _patch_plotly()
        if name == "reportlab.pdfgen.canvas" or (name == "reportlab.pdfgen" and "canvas" in (fromlist or ())):
            _patch_reportlab()
        if name == "docx":
            _patch_docx()
        if name == "pptx":
            _patch_pptx()
    except Exception:
        pass
    return mod


def ensure_auto_artifact_hooks(workdir: str):
    global _AUTO_HOOKS_INSTALLED, _WORKDIR
    _WORKDIR = workdir
    _ensure_matplotlib_env(workdir)

    if not _AUTO_HOOKS_INSTALLED:
        builtins.__import__ = _croki_import
        _AUTO_HOOKS_INSTALLED = True

    if "matplotlib.pyplot" in sys.modules:
        _patch_matplotlib()
    if "PIL.Image" in sys.modules or "PIL" in sys.modules:
        _patch_pil()
    if "plotly" in sys.modules:
        _patch_plotly()
    if "reportlab.pdfgen.canvas" in sys.modules:
        _patch_reportlab()
    if "docx" in sys.modules:
        _patch_docx()
    if "pptx" in sys.modules:
        _patch_pptx()
