import json
import os
import signal
import subprocess
import sys
import time
import unittest
from pathlib import Path


_ROOT = Path(__file__).resolve().parents[1]
_SERVER = _ROOT / "server.py"
_CLIENT = _ROOT / "sandbox_client.py"
_OUTPUTS = Path(__file__).resolve().parent / "outputs"


def _wait_for_socket(path: Path, timeout_seconds: float) -> None:
    start = time.time()
    while time.time() - start < timeout_seconds:
        if path.exists():
            return
        time.sleep(0.05)
    raise RuntimeError(f"socket did not appear: {path}")


class Sandbox:
    def __init__(self, workdir: Path):
        self.workdir = workdir
        self.socket_path = workdir / "sandbox.sock"
        self.proc: subprocess.Popen | None = None

    def start(self) -> None:
        self.workdir.mkdir(parents=True, exist_ok=True)
        try:
            self.socket_path.unlink()
        except FileNotFoundError:
            pass

        env = os.environ.copy()
        env["SANDBOX_WORKDIR"] = str(self.workdir)
        env["SANDBOX_SOCKET"] = str(self.socket_path)
        env["SANDBOX_TIMEOUT_SECONDS"] = "30"

        self.proc = subprocess.Popen(
            [sys.executable, str(_SERVER)],
            cwd=str(_ROOT),
            env=env,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        _wait_for_socket(self.socket_path, timeout_seconds=5.0)

    def stop(self) -> None:
        if not self.proc:
            return
        if self.proc.poll() is None:
            self.proc.send_signal(signal.SIGTERM)
            try:
                self.proc.wait(timeout=3)
            except subprocess.TimeoutExpired:
                self.proc.kill()
                self.proc.wait(timeout=3)
        self.proc = None

    def exec(self, code: str, timeout_seconds: float = 5.0) -> dict:
        env = os.environ.copy()
        env["SANDBOX_SOCKET"] = str(self.socket_path)
        raw = json.dumps({"code": code}, ensure_ascii=False) + "\n"
        try:
            p = subprocess.run(
                [sys.executable, str(_CLIENT)],
                input=raw.encode("utf-8"),
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                env=env,
                timeout=timeout_seconds,
                check=True,
            )
        except subprocess.TimeoutExpired:
            self.stop()
            self.start()
            return {"stdout": "", "stderr": "", "error": "context deadline exceeded", "changed_files": []}
        out = p.stdout.decode("utf-8", errors="replace").strip()
        return json.loads(out)


class TestArtifacts(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _OUTPUTS.mkdir(parents=True, exist_ok=True)
        cls.sandbox = Sandbox(_OUTPUTS)
        cls.sandbox.start()

    @classmethod
    def tearDownClass(cls) -> None:
        cls.sandbox.stop()

    def _assert_has_extension(self, resp: dict, ext: str) -> None:
        self.assertEqual(resp.get("error", ""), "")
        changed = resp.get("changed_files") or []
        self.assertIsInstance(changed, list)
        matches = [Path(p) for p in changed if str(p).endswith(ext)]
        self.assertTrue(matches, f"no {ext} in changed_files: {changed}")
        for p in matches:
            self.assertTrue(p.exists(), f"missing file: {p}")
            self.assertGreater(p.stat().st_size, 0, f"empty file: {p}")

    def test_matplotlib_show_creates_png(self) -> None:
        resp = self.sandbox.exec(
            "import matplotlib.pyplot as plt\n"
            "plt.plot([1,2,3],[1,4,9])\n"
            "plt.axis('off')\n"
            "plt.show()\n",
            timeout_seconds=15.0,
        )
        self._assert_has_extension(resp, ".png")

    def test_matplotlib_savefig_reports_png(self) -> None:
        out = _OUTPUTS / "instagram_logo.png"
        try:
            out.unlink()
        except FileNotFoundError:
            pass

        resp = self.sandbox.exec(
            "import matplotlib.pyplot as plt\n"
            "import numpy as np\n"
            "fig, ax = plt.subplots(figsize=(3, 3))\n"
            "ax.axis('off')\n"
            "z = np.linspace(0, 1, 64).reshape(8, 8)\n"
            "ax.imshow(z, interpolation='bicubic', cmap='plasma', origin='lower')\n"
            "plt.savefig('instagram_logo.png', bbox_inches='tight', pad_inches=0.1, dpi=100)\n"
            "plt.close()\n",
            timeout_seconds=15.0,
        )
        self._assert_has_extension(resp, ".png")
        self.assertTrue(out.exists(), f"missing file: {out}")

    def test_invalid_matplotlib_colormap_returns_error_and_sandbox_continues(self) -> None:
        resp = self.sandbox.exec(
            "import matplotlib.pyplot as plt\n"
            "import numpy as np\n"
            "fig, ax = plt.subplots()\n"
            "ax.imshow(np.zeros((2,2)), cmap='instagram_gradient')\n"
            "plt.close(fig)\n",
            timeout_seconds=15.0,
        )
        self.assertIn("ValueError", resp.get("error", ""))
        self.assertIn("instagram_gradient", resp.get("error", ""))

        resp2 = self.sandbox.exec("print('ok')\n")
        self.assertEqual(resp2.get("error", ""), "")
        self.assertIn("ok", resp2.get("stdout", ""))

    def test_pil_show_creates_png(self) -> None:
        resp = self.sandbox.exec(
            "from PIL import Image\n"
            "img = Image.new('RGB', (32, 32), (255, 0, 0))\n"
            "img.show()\n"
        )
        self._assert_has_extension(resp, ".png")

    def test_plotly_show_creates_html(self) -> None:
        resp = self.sandbox.exec(
            "import plotly.graph_objects as go\n"
            "fig = go.Figure(data=go.Scatter(y=[1, 3, 2]))\n"
            "fig.show()\n"
        )
        self._assert_has_extension(resp, ".html")

    def test_pdf_is_reported(self) -> None:
        resp = self.sandbox.exec(
            "from reportlab.pdfgen import canvas\n"
            "c = canvas.Canvas('test.pdf')\n"
            "c.drawString(72, 720, 'hello')\n"
            "c.save()\n"
        )
        self._assert_has_extension(resp, ".pdf")

    def test_docx_is_reported(self) -> None:
        resp = self.sandbox.exec(
            "from docx import Document\n"
            "d = Document()\n"
            "d.add_paragraph('hello')\n"
            "d.save('test.docx')\n"
        )
        self._assert_has_extension(resp, ".docx")

    def test_doc_is_reported(self) -> None:
        resp = self.sandbox.exec(
            "xml = \"\"\"<?xml version='1.0' encoding='UTF-8'?>\n"
            "<w:wordDocument xmlns:w='http://schemas.microsoft.com/office/word/2003/wordml'>\n"
            "  <w:body>\n"
            "    <w:p><w:r><w:t>hello</w:t></w:r></w:p>\n"
            "  </w:body>\n"
            "</w:wordDocument>\n"
            "\"\"\"\n"
            "open('test.doc', 'w', encoding='utf-8').write(xml)\n"
        )
        self._assert_has_extension(resp, ".doc")

    def test_pptx_is_reported(self) -> None:
        resp = self.sandbox.exec(
            "from pptx import Presentation\n"
            "p = Presentation()\n"
            "p.slides.add_slide(p.slide_layouts[5])\n"
            "p.save('test.pptx')\n"
        )
        self._assert_has_extension(resp, ".pptx")
