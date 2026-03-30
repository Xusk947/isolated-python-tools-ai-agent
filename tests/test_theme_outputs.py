import unittest
import zipfile
from pathlib import Path

from PIL import Image, ImageChops

from helpers import assert_has_extension
from test_artifacts import Sandbox


_OUT_DIR = Path(__file__).resolve().parent / "outputs" / "theme"


UZBEKISTAN_TEXT = (
    "Uzbekistan lies in the heart of Central Asia and has long been a crossroads of empires, "
    "trade, and cultures along the Silk Road. Over centuries it was ruled by Persian, Arab, "
    "Mongol–Timurid, and later Russian and Soviet powers before becoming an independent republic "
    "in 1991. Since independence, Uzbekistan has focused on building its own state institutions, "
    "economy, and foreign‑policy path while balancing its rich Islamic and Soviet‑era heritage."
)


class TestThemeOutputs(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        _OUT_DIR.mkdir(parents=True, exist_ok=True)
        cls.sandbox = Sandbox(_OUT_DIR)
        cls.sandbox.start()

    @classmethod
    def tearDownClass(cls) -> None:
        cls.sandbox.stop()

    def test_matplotlib_population_plot_is_not_blank(self) -> None:
        resp = self.sandbox.exec(
            "import matplotlib.pyplot as plt\n"
            "years = list(range(2000, 2026))\n"
            "Kazakhstan = [14883626,14953571,15013871,15075901,15143320,15211800,15275600,15331300,15379000,15418000,15443000,15449000,15431000,15391000,15331000,16180000,16780000,17330000,17790000,18220000,18860000,18930000,19130000,20030000,20590000,20480000]\n"
            "Uzbekistan = [24650400,25120000,25570000,26000000,26410000,26780000,27110000,27400000,27660000,27890000,28100000,28290000,28470000,28640000,28820000,30110000,31010000,31830000,32600000,33220000,33590000,34240000,34940000,35650000,36360000,37050000]\n"
            "Kyrgyzstan = [5025140,5118000,5206000,5290000,5368000,5440000,5502000,5552000,5591000,5619000,5635000,5638000,5631000,5614000,5591000,6010000,6300000,6480000,6600000,6680000,6730000,6850000,6980000,7100000,7220000,7290000]\n"
            "plt.figure(figsize=(9, 4.5))\n"
            "plt.plot(years, Kazakhstan, label='Kazakhstan')\n"
            "plt.plot(years, Uzbekistan, label='Uzbekistan')\n"
            "plt.plot(years, Kyrgyzstan, label='Kyrgyzstan')\n"
            "plt.title('Population (2000–2025)')\n"
            "plt.xlabel('Year')\n"
            "plt.ylabel('Population')\n"
            "plt.legend()\n"
            "plt.tight_layout()\n"
            "plt.show()\n",
            timeout_seconds=15.0,
        )
        p = assert_has_extension(resp, ".png")[-1]
        img = Image.open(p).convert("RGB")
        bg = img.getpixel((0, 0))
        diff = ImageChops.difference(img, Image.new(img.mode, img.size, bg))
        self.assertIsNotNone(diff.getbbox(), "plot image seems blank")

    def test_plotly_population_plot_includes_theme_fonts(self) -> None:
        resp = self.sandbox.exec(
            "import plotly.graph_objects as go\n"
            "data = {\n"
            "  'Kazakhstan': {'2000': 14883626, '2020': 18855767, '2021': 18927000, '2022': 19131779, '2023': 20033546, '2024': 20592571, '2025': 20478879},\n"
            "  'Uzbekistan': {'2000': 24650400, '2020': 33586372, '2021': 34243696, '2022': 34938955, '2023': 35652307, '2024': 36361859, '2025': 37053438},\n"
            "  'Kyrgyzstan': {'2000': 5025140, '2020': 6726596, '2021': 6851224, '2022': 6975220, '2023': 7099750, '2024': 7221868, '2025': 7295030},\n"
            "}\n"
            "fig = go.Figure()\n"
            "for country, pts in data.items():\n"
            "  xs = sorted(int(k) for k in pts.keys())\n"
            "  ys = [pts[str(x)] for x in xs]\n"
            "  fig.add_trace(go.Scatter(x=xs, y=ys, mode='lines+markers', name=country))\n"
            "fig.update_layout(title='Population (Selected Years)', xaxis_title='Year', yaxis_title='Population')\n"
            "fig.show()\n",
            timeout_seconds=15.0,
        )
        p = assert_has_extension(resp, ".html")[-1]
        raw = p.read_text(encoding="utf-8", errors="replace")
        self.assertIn("Kazakhstan", raw)
        self.assertIn("Lato", raw)
        self.assertIn("Unbounded", raw)

    def test_docx_has_theme_fonts_and_color(self) -> None:
        resp = self.sandbox.exec(
            "from docx import Document\n"
            "doc = Document()\n"
            "doc.add_heading('Uzbekistan', level=1)\n"
            f"doc.add_paragraph({UZBEKISTAN_TEXT!r})\n"
            "doc.save('uzbekistan.docx')\n",
            timeout_seconds=15.0,
        )
        p = assert_has_extension(resp, ".docx")[-1]
        with zipfile.ZipFile(p, "r") as zf:
            styles = zf.read("word/styles.xml").decode("utf-8", errors="replace")
        self.assertIn("Lato", styles)
        self.assertIn("Unbounded", styles)
        self.assertIn("111827", styles)

    def test_pdf_references_lato(self) -> None:
        resp = self.sandbox.exec(
            "from reportlab.pdfgen import canvas\n"
            "from reportlab.lib.pagesizes import letter\n"
            "c = canvas.Canvas('uzbekistan.pdf', pagesize=letter)\n"
            "c.drawString(72, 740, 'Uzbekistan')\n"
            f"c.drawString(72, 720, {UZBEKISTAN_TEXT[:110]!r})\n"
            "c.showPage()\n"
            "c.save()\n",
            timeout_seconds=15.0,
        )
        p = assert_has_extension(resp, ".pdf")[-1]
        self.assertIn(b"Lato", p.read_bytes())

    def test_pptx_has_theme_fonts(self) -> None:
        half = len(UZBEKISTAN_TEXT) // 2
        part1 = UZBEKISTAN_TEXT[:half]
        part2 = UZBEKISTAN_TEXT[half:]
        resp = self.sandbox.exec(
            "from pptx import Presentation\n"
            "prs = Presentation()\n"
            "s1 = prs.slides.add_slide(prs.slide_layouts[0])\n"
            "s1.shapes.title.text = 'Uzbekistan'\n"
            f"s1.placeholders[1].text = {part1!r}\n"
            "s2 = prs.slides.add_slide(prs.slide_layouts[1])\n"
            "s2.shapes.title.text = 'History'\n"
            f"s2.placeholders[1].text = {part2!r}\n"
            "prs.save('uzbekistan.pptx')\n",
            timeout_seconds=15.0,
        )
        p = assert_has_extension(resp, ".pptx")[-1]
        with zipfile.ZipFile(p, "r") as zf:
            slides = [zf.read(n).decode("utf-8", errors="replace") for n in zf.namelist() if n.startswith("ppt/slides/slide")]
        merged = "\n".join(slides)
        self.assertIn('typeface="Unbounded"', merged)
        self.assertIn('typeface="Lato"', merged)
